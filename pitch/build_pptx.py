"""
Generate pitch/ParkPulse_Deck.pptx — a dark-themed submission deck (judge- & AI-readable, self-contained).
Run:  pip install python-pptx  &&  python pitch/build_pptx.py
Mirrors pitch/PITCH_DECK.md (keep them in sync).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0x0B, 0x0E, 0x14)
PANEL  = RGBColor(0x12, 0x16, 0x1F)
FG     = RGBColor(0xE6, 0xE9, 0xEF)
MUTED  = RGBColor(0x9A, 0xA3, 0xB2)
ACCENT = RGBColor(0x4C, 0x8B, 0xF5)
RED    = RGBColor(0xE2, 0x35, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    # accent bar top-left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.7), Inches(0.55), Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    """runs: list of (text, size, color, bold)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = "Segoe UI"
    return tb


def title_slide():
    s = prs.slides.add_slide(BLANK); _bg(s)
    _text(s, Inches(0.9), Inches(2.2), Inches(11.8), Inches(3.2),
          [("ParkPulse", 60, FG, True),
           ("From 298,000 parking tickets to where to stand tomorrow.", 26, ACCENT, False),
           ("Enforcement intelligence for the Bengaluru Traffic Police — built on their own challan data.", 16, MUTED, False),
           ("", 8, MUTED, False),
           ("Theme 1 — Poor Visibility on Parking-Induced Congestion", 16, MUTED, False),
           ("Gridlock Hackathon 2.0  ·  TeamX", 16, MUTED, False)])
    return s


def content_slide(kicker, title, bullets, footer=None):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _text(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.4), [(kicker, 13, ACCENT, True)])
    _text(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.0), [(title, 32, FG, True)])
    runs = [("•  " + b, 18, FG, False) for b in bullets]
    _text(s, Inches(1.0), Inches(2.7), Inches(11.4), Inches(3.8), runs, space=11)
    if footer:
        _text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.6), [(footer, 13, MUTED, False)])
    return s


def stat_slide(kicker, title, big, big_label, bullets):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _text(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.4), [(kicker, 13, ACCENT, True)])
    _text(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.0), [(title, 32, FG, True)])
    # big number
    _text(s, Inches(0.9), Inches(2.8), Inches(4.6), Inches(2.6),
          [(big, 110, ACCENT, True), (big_label, 16, MUTED, False)], anchor=MSO_ANCHOR.MIDDLE)
    runs = [("•  " + b, 18, FG, False) for b in bullets]
    _text(s, Inches(5.9), Inches(2.9), Inches(6.6), Inches(3.4), runs, space=11)
    return s


# 1
title_slide()
# 2
content_slide("The problem · Relevance", "Enforcement is flying blind", [
    "Illegal parking near markets, metros & hubs blocks live lanes and junctions — the city's most visible congestion source.",
    "Enforcement is reactive & patrol-based: a few teams, a whole division, sent on instinct.",
    "No city-wide heatmap · no prioritization · no way to schedule scarce teams by where/when parking actually hurts."])
# 3
content_slide("The data · Feasibility", "We didn't invent data — we sharpened theirs", [
    "298,445 real BTP violation records · 151 days · 802 zones · zero missing fields.",
    "Wildly concentrated: the top 1% of locations hold 33% of all violations.",
    "15% of vehicles cause 34% of violations — a chronic-offender tail.",
    "That concentration is the opportunity: a few well-placed teams cover a disproportionate share."],
    footer="Their own challan data — no new sensors, nothing invented.")
# 4
content_slide("The honest insight · Trust", "We read the data for what it is — enforcement, not demand", [
    "It shows where officers caught parking, not a god's-eye view of where parking is worst.",
    "Morning-heavy: ~93% of tickets are written before 1 PM; <0.3% in the 5–9 PM evening peak.",
    "So we make two honest claims, not one inflated one: optimize the effort you already spend,",
    "and flag the evening blind spot you're missing — instead of pretending to predict the whole city."])
# 5
content_slide("The solution · Innovation", "Not a dashboard — a closed decision loop", [
    "Detect → Score → Forecast → Deploy → Target → (outcomes feed back).",
    "It ends in an action a supervisor can brief at tomorrow's roll-call, not a chart to stare at.",
    "Every stage is transparent and explainable — no black boxes to take on faith."])
# 6
content_slide("Detect & Score · Technical", "One transparent number ranks every zone", [
    "City-wide 3-D hotspot map over geohash hexbins (gh6 ≈ 1.2 km, gh7 ≈ 150 m).",
    "Congestion Impact Score 0–100 = violations × avg severity × flow-criticality.",
    "Junction-blocking and main-road violations weigh most — 100 cars on a footpath ≠ 100 blocking a junction.",
    "Fully explainable: anyone can see why a zone scores what it does — no learned weights to defend."])
# 7
content_slide("Forecast · Innovation + credibility", "A forecast we tested on the future, not the past", [
    "Predicts violations per zone × weekday × hour via an empirical-Bayes shrinkage model (sparse cells stay stable).",
    "Honest held-out backtest: r = 0.70, MAE 2.01 — trained on early weeks, scored on weeks it never saw.",
    "Benchmarked LightGBM: it scored lower (0.69) and we couldn't explain it → kept the interpretable model."])
# 8
stat_slide("Proven impact · Real-world impact", "Proven, not claimed", "38%", "captured · 10 teams · 31 unseen days", [
    "Replayed a 10-team plan across 31 held-out days the model never saw.",
    "Those teams would have been sitting on ~38% of the violations actually logged…",
    "…vs ~1.3% spread evenly (and well above a static 'usual hotspots' plan too).",
    "A validated efficiency gain on unseen data — most teams have zero validation."])
# 9
content_slide("Deploy & Target · Impact + Trust", "The deliverable an officer uses — auditable both ways", [
    "One click → a spaced patrol plan (teams ≥600 m apart), shared by WhatsApp / print / CSV at roll-call.",
    "Why a team IS there: tap it for forecast load · junction/main-road impact · 3-week trend.",
    "Why a strong zone ISN'T (new): an 'Also considered' panel shows what just missed and why — auditable, not a black box.",
    "Target chronic offenders (15% → 34%): ready owner-notice / escalated-penalty / tow-priority CSV lists."])
# 10
content_slide("Ask ParkPulse · Innovation", "A real AI agent over the live models — not a chatbot", [
    "Gemini function-calling runs the actual forecaster & optimizer; every number is grounded in a tool result.",
    "Multilingual (English / Hindi / Kannada) + voice input + conversation memory.",
    "An officer just asks or says: 'Plan 6 teams for Friday evening near KR Market.'"])
# 11
content_slide("By design · Integrity", "Live, explainable & self-improving", [
    "Event-aware: auto-flags unusual days (festivals / match days) above the weekday norm.",
    "Full-day planner: morning → afternoon → evening at once (361 → 36 → 2) — see the evening blind spot.",
    "Data flywheel: a live ingest endpoint cleans new challans + rebuilds models in seconds; officers log real outcomes.",
    "Integrity: we never modify the fixed dataset — ingest is live-data architecture, outcome logs in-memory only."])
# 12
content_slide("Architecture & tech · Technical", "One engine, two services — small, fast, explainable", [
    "One engine (core.py, pandas/numpy): EB-shrinkage forecaster · impact score · spaced greedy optimizer · honest backtest · simulator.",
    "The product: Next.js 16 / React 19 / deck.gl frontend (Vercel) + FastAPI backend (Render) — both call the one engine.",
    "Gemini co-pilot runs server-side over the real models — no GPU, no training pipeline to babysit, no black box."])
# 13
content_slide("Feasibility · rollout", "Deployable today — on data BTP already collects", [
    "Runs on a laptop · only the existing challan data · no new sensors, no procurement.",
    "Live as a web app today; the full-stack product is deployment-ready (FastAPI on Render + Next.js on Vercel).",
    "The geohash engine generalizes to any city or violation type — not parking-specific.",
    "Next step for BTP: a nightly challan feed + a few evening pilot deployments to start the flywheel."])
# 14
content_slide("Close", "Why ParkPulse wins", [
    "Feasible — working software today, on data BTP already has, no new hardware.",
    "Relevant — their data, their problem, read honestly (enforcement ≠ demand; evening blind spot named).",
    "Innovative — a closed decision loop + an agentic voice co-pilot + transparency both ways.",
    "Real impact — a plan you brief tomorrow, a 38% validated efficiency gain, offender targeting.",
    "Turning what already happened into where to stand tomorrow."],
    footer="Live app: Next.js on Vercel + FastAPI on Render  ·  [Vercel URL]")

out = Path(__file__).resolve().parent / "ParkPulse_Deck.pptx"
prs.save(str(out))
print("saved:", out)
