"""
Generate pitch/ParkPulse_Submission_Deck.pptx — the official Gridlock-template deck (10 slides, judge/AI-readable).
Run:  pip install python-pptx  &&  python pitch/build_submission_pptx.py
Mirrors pitch/SUBMISSION_DECK.md (keep them in sync). Replace [links] before submitting.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0x0B, 0x0E, 0x14)
FG     = RGBColor(0xE6, 0xE9, 0xEF)
MUTED  = RGBColor(0x9A, 0xA3, 0xB2)
ACCENT = RGBColor(0x4C, 0x8B, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.7), Inches(0.55), Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = "Segoe UI"
    return tb


def title_slide(kicker, big, sub, bullets, footer):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.4), [(kicker, 13, ACCENT, True)])
    _text(s, Inches(0.9), Inches(1.5), Inches(11.6), Inches(2.0),
          [(big, 50, FG, True), (sub, 22, ACCENT, False)])
    runs = [("•  " + b, 18, FG, False) for b in bullets]
    _text(s, Inches(1.0), Inches(3.9), Inches(11.4), Inches(2.6), runs, space=11)
    _text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.6), [(footer, 14, MUTED, False)])
    return s


def content_slide(kicker, title, bullets, footer=None):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _text(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.4), [(kicker, 13, ACCENT, True)])
    _text(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.0), [(title, 31, FG, True)])
    runs = [("•  " + b, 18, FG, False) for b in bullets]
    _text(s, Inches(1.0), Inches(2.7), Inches(11.4), Inches(3.9), runs, space=10)
    if footer:
        _text(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.5), [(footer, 13, MUTED, False)])
    return s


def stat_slide(kicker, title, big, big_label, bullets):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _text(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.4), [(kicker, 13, ACCENT, True)])
    _text(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.0), [(title, 31, FG, True)])
    _text(s, Inches(0.9), Inches(2.8), Inches(4.6), Inches(2.6),
          [(big, 110, ACCENT, True), (big_label, 16, MUTED, False)], anchor=MSO_ANCHOR.MIDDLE)
    runs = [("•  " + b, 18, FG, False) for b in bullets]
    _text(s, Inches(5.9), Inches(2.9), Inches(6.6), Inches(3.4), runs, space=11)
    return s


# 1 — Context & Scope
title_slide("Slide 1 · Context & Scope",
            "ParkPulse",
            "From 298,000 parking tickets to where to stand tomorrow.",
            ["Scope: illegal on-street parking near markets, metros & commercial hubs — vehicles in live lanes and at junctions.",
             "Goal: a data-driven enforcement system on the BTP's own challan data — which parking hurts, when it recurs, where to send teams.",
             "Theme 1 — Poor Visibility on Parking-Induced Congestion  ·  Gridlock Hackathon 2.0  ·  TeamX"],
            "Bengaluru Traffic Police · built on 298,445 real violation records")
# 2 — Problem Description
content_slide("Slide 2 · Problem Description", "Enforcement is reactive and blind", [
    "A duty officer has a few teams and a whole division, and decides where to send them on instinct — no city-wide picture.",
    "Existing methods fall short: manual patrols don't prioritize; cameras are capital-heavy and don't say where to deploy people;",
    "static 'known hotspot' lists ignore WHEN violations happen and go stale.",
    "The data to fix it already exists (challan records) but sits unused — no heatmap, no ranking, no schedule."])
# 3 — Analysis Part 1: The "Why"
content_slide("Slide 3 · Analysis & Research, Part 1 — The “Why”", "The problem is real, and extremely concentrated", [
    "298,445 real BTP violation records · 151 days · 802 geohash zones · zero missing fields — their own data, cleaned.",
    "Root cause is spatial concentration: the top 1% of locations hold 33% of all violations.",
    "Plus a chronic-offender tail: 15% of vehicles cause 34% of violations.",
    "A few well-placed teams can therefore cover a disproportionate share of the problem."])
# 4 — Analysis Part 2: Impact
content_slide("Slide 4 · Analysis & Research, Part 2 — Impact", "It hits commuters, business & delivery — and we read it honestly", [
    "Illegal parking in live lanes and at junctions slows commuters, blocks bus stops/main roads, chokes last-mile delivery.",
    "Honest insight: this is ENFORCEMENT data, not demand — ~93% of tickets are written before 1 PM, <0.3% in the 5–9 PM peak.",
    "So we make two honest claims: optimize the effort already spent, AND flag the evening blind spot — no over-claiming."])
# 5 — Solution Part 1: Concept
content_slide("Slide 5 · Proposed Solution, Part 1 — The Concept", "Not a dashboard — a closed decision loop", [
    "Detect → Score → Forecast → Deploy → Target → (outcomes feed back). It ends in a plan a supervisor briefs tomorrow.",
    "Impact Score (0–100) = violations × avg severity × flow-criticality (junction/main-road weight).",
    "Forecaster: empirical-Bayes shrinkage of violations per zone × weekday × hour (sparse cells stay stable).",
    "Optimizer: greedy maximum-coverage with a ≥600 m spacing constraint. Every stage is explainable — no black box."])
# 6 — Solution Part 2: Architecture & Tech Stack
content_slide("Slide 6 · Proposed Solution, Part 2 — Architecture & Tech Stack", "One engine, two services — small, fast, auditable", [
    "One engine (core.py, pandas/numpy): impact score, EB-shrinkage forecaster, spaced optimizer, honest backtest, simulator.",
    "The product: a Next.js 16 / React 19 / deck.gl frontend on Vercel + a FastAPI backend on Render — both call the one engine.",
    "AI co-pilot: Google Gemini automatic function-calling runs the real tools server-side (key never leaves the server).",
    "No GPU, no training pipeline, no IoT hardware — the whole model rebuilds in seconds."])
# 7 — Solution Part 3: Implementation
content_slide("Slide 7 · Proposed Solution, Part 3 — Implementation", "How it's executed — and what sets it apart", [
    "Pick weekday + shift + team count → a spaced patrol plan, shared by WhatsApp / print / CSV at roll-call.",
    "Auditable both ways: tap a team for WHY it's placed; an 'Also considered' panel shows why a strong zone ISN'T.",
    "Full-day planner exposes the evening blind spot (361 → 36 → 2) and nudges reallocation into the gap.",
    "Co-pilot: ask or SPEAK in English / Hindi / Kannada. Target chronic offenders (15% → 34%) with ready CSV lists."])
# 8 — Validation & Proven Impact
stat_slide("Slide 8 · Validation & Proven Impact", "Proven on data the model never saw", "38%", "captured · 10 teams · 31 unseen days", [
    "Forecast validated the hard way: r = 0.70, MAE 2.01 on the unseen future; beat a LightGBM model (0.69).",
    "Counterfactual: a 10-team plan replayed on 31 held-out days would have sat on ~38% of the violations actually logged…",
    "…vs ~1.3% spread evenly (and well above a static 'usual hotspots' plan).",
    "A validated efficiency gain on unseen data — most teams have zero validation."])
# 9 — Feasibility, Rollout & Integrity
content_slide("Slide 9 · Feasibility, Rollout & Integrity", "Deployable today — and we never touched the fixed data", [
    "Runs on a laptop · only existing challan data · no new sensors or procurement. Live as a web app today.",
    "Self-improving: a live ingest endpoint rebuilds models in seconds (in-memory); outcome logging closes the loop.",
    "Integrity: we NEVER modify the fixed competition dataset — ingest is architecture for live data; logs are in-memory only.",
    "Roadmap: ① Retrospective (today) → ② Nightly challan feed + evening pilots → ③ Fuse the live ASTraM / camera feed."])
# 10 — Supporting Links & Documentation
content_slide("Slide 10 · Supporting Links & Documentation", "Everything to verify it runs", [
    "Live web app (Next.js on Vercel): [Vercel URL]",
    "Backend API (FastAPI on Render): [Render URL]/health",
    "Code repository: [GitHub repo URL]   ·   Demo video: [link — ~3–5 min walkthrough]",
    "Run locally: python fullstack/backend/main.py  ·  cd fullstack/frontend && npm install && npm run dev",
    "Team: TeamX — Shubham, Satyam, Palash.  ParkPulse: turning what already happened into where to stand tomorrow."])

out = Path(__file__).resolve().parent / "ParkPulse_Submission_Deck.pptx"
prs.save(str(out))
print("saved:", out)
